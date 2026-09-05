# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[2]
SOFTWAREBAR = ROOT / "softwarebar"
DOCS = SOFTWAREBAR / "docs"
ASSETS = SOFTWAREBAR / "frontend-vue" / "public" / "assets"
OUT = DOCS / "产品方案介绍PPT_高分版.pptx"


W, H = Cm(33.867), Cm(19.05)

COLORS = {
    "ink": RGBColor(23, 35, 50),
    "muted": RGBColor(91, 107, 128),
    "light": RGBColor(246, 248, 251),
    "line": RGBColor(218, 226, 237),
    "blue": RGBColor(37, 99, 235),
    "teal": RGBColor(13, 148, 136),
    "green": RGBColor(22, 163, 74),
    "amber": RGBColor(217, 119, 6),
    "red": RGBColor(220, 38, 38),
    "purple": RGBColor(124, 58, 237),
    "white": RGBColor(255, 255, 255),
    "dark": RGBColor(15, 23, 42),
}


def add_bg(slide, color="light"):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS[color]


def set_run(run, size=18, bold=False, color="ink", font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color] if isinstance(color, str) else color


def add_text(slide, text, x, y, w, h, size=18, bold=False, color="ink", align=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 1.4, 0.62, 23.8, 0.95, size=24, bold=True)
    if subtitle:
        add_text(slide, subtitle, 1.45, 1.48, 24.6, 0.5, size=9.5, color="muted")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.4), Cm(2.08), Cm(31.0), Cm(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_card(slide, x, y, w, h, title=None, body=None, accent="blue", fill="white"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[fill]
    shp.line.color.rgb = COLORS["line"]
    shp.line.width = Pt(0.8)
    shp.adjustments[0] = 0.08
    if title:
        add_text(slide, title, x + 0.38, y + 0.28, w - 0.76, 0.42, size=12, bold=True, color=accent)
    if body:
        add_text(slide, body, x + 0.38, y + (0.86 if title else 0.36), w - 0.76, h - 0.95, size=9.6, color="ink")
    return shp


def add_bullets(slide, items, x, y, w, h, size=11, color="ink", gap=0.2):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap * 10)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
    return box


def add_metric(slide, label, value, x, y, w, accent="blue"):
    add_card(slide, x, y, w, 2.15, fill="white")
    add_text(slide, value, x + 0.28, y + 0.34, w - 0.56, 0.7, size=24, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.28, y + 1.18, w - 0.56, 0.48, size=9.5, color="muted", align=PP_ALIGN.CENTER)


def add_pill(slide, text, x, y, w, accent="blue"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.64))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[accent]
    shp.line.fill.background()
    shp.adjustments[0] = 0.45
    add_text(slide, text, x, y + 0.11, w, 0.35, size=8.6, bold=True, color="white", align=PP_ALIGN.CENTER)


def add_image(slide, rel, x, y, w, h):
    path = ASSETS / rel
    if path.exists():
        slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    else:
        add_card(slide, x, y, w, h, "视觉素材", f"缺少图片：{rel}", accent="amber")


def cover():
    slide = prs.slides.add_slide(blank)
    add_bg(slide, "dark")
    add_image(slide, "scenic/landing-hero.png", 16.0, 0.0, 17.9, 19.05)
    add_text(slide, "A5 多模态 AI 智能景区导览助手", 1.55, 2.45, 15.8, 1.2, size=32, bold=True, color="white")
    add_text(slide, "灵灵导游：数字人讲解 + 本地知识库 RAG + 个性化路线 + 运营后台", 1.62, 3.82, 15.0, 0.7, size=14, color="white")
    add_card(slide, 1.62, 5.15, 12.4, 3.1, fill="white", accent="blue")
    add_text(slide, "面向景区的可落地智能导览系统", 2.1, 5.68, 11.3, 0.48, size=15, bold=True, color="ink")
    add_bullets(slide, [
        "游客端：语音/文本/情绪互动，景点问答、讲解、路线推荐",
        "数字人端：Wav2Lip 口型驱动，支持 WebRTC 播放与音色切换",
        "管理端：知识库、形象、反馈与客流行为数据看板",
    ], 2.1, 6.35, 11.2, 1.55, size=9.5)
    add_pill(slide, "功能完整", 1.65, 9.05, 3.1, "blue")
    add_pill(slide, "技术可信", 5.0, 9.05, 3.1, "teal")
    add_pill(slide, "可本地运行", 8.35, 9.05, 3.6, "green")
    add_text(slide, "参赛作品产品方案介绍", 1.62, 17.35, 8.8, 0.38, size=9, color="white")


def scoring():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "围绕赛题评分项组织方案", "PPT 内容直接对齐功能、技术、落地和文档四类评价指标")
    items = [
        ("40%", "功能完整性", "游客交互、景点讲解、路线推荐、数字人、后台运营形成闭环", "blue"),
        ("30%", "技术与创新", "DeepSeek + Qwen-VL + RAG + Wav2Lip + 数据看板的组合式架构", "purple"),
        ("20%", "行业适用性", "基于景区真实业务流程，兼顾游客体验和景区运营管理", "teal"),
        ("10%", "文档与展示", "部署说明、测试结果、演示脚本、PPT 与视频脚本可直接交付", "green"),
    ]
    for i, (pct, title, desc, accent) in enumerate(items):
        x = 1.5 + i * 7.8
        add_card(slide, x, 3.1, 6.9, 8.3, fill="white")
        add_text(slide, pct, x + 0.35, 3.6, 6.2, 0.85, size=30, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.4, 4.78, 6.1, 0.45, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.65, 5.65, 5.6, 1.3, size=10.2, color="muted", align=PP_ALIGN.CENTER)
        add_pill(slide, ["必须项覆盖", "模型组合", "业务闭环", "可验收"][i], x + 1.55, 9.5, 3.8, accent)
    add_text(slide, "答辩主线：游客侧体验可感知、技术侧指标可验证、景区侧价值可运营。", 2.1, 13.25, 29.4, 0.7, size=16, bold=True, color="ink", align=PP_ALIGN.CENTER)


def pain_points():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "行业痛点与产品定位", "传统导览往往停留在单点信息展示，本项目把“问、讲、游、管”连成闭环")
    left = [
        ("讲解不够个性化", "固定语音或图文说明无法适配亲子、研学、文化深度游等差异化需求。"),
        ("咨询响应不稳定", "高峰期人工问询压力大，游客重复问题多，服务质量受排队影响。"),
        ("路线规划粗粒度", "路线推荐缺少时长、偏好、实时位置等条件，无法动态引导。"),
        ("运营数据难沉淀", "游客反馈和行为数据分散，难以反哺内容优化与服务管理。"),
    ]
    for i, (t, b) in enumerate(left):
        add_card(slide, 1.55, 2.8 + i * 2.55, 13.8, 2.0, t, b, accent=["red", "amber", "purple", "teal"][i])
    add_image(slide, "scenic/photos/lingshan-grand-buddha.jpg", 17.0, 2.8, 14.0, 7.6)
    add_card(slide, 17.0, 11.0, 14.0, 4.25, "产品定位", "以“灵灵导游”为统一入口，用多模态大模型理解游客问题，用本地知识库保障景区事实准确，用数字人提升讲解沉浸感，并把互动、反馈、行为数据沉淀到管理后台。", accent="blue")


def loop():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "产品闭环：游客体验与景区运营同步提升")
    steps = [
        ("1", "游客提问", "文本/语音输入，识别情绪和偏好"),
        ("2", "知识检索", "景区资料、官方文档、行为数据联合召回"),
        ("3", "智能生成", "大模型组织讲解、问答和路线推荐"),
        ("4", "数字人播报", "口型驱动与音色合成，形成沉浸式讲解"),
        ("5", "数据回流", "反馈、满意度、热点景点进入运营看板"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 1.45 + i * 6.25
        add_card(slide, x, 4.25, 5.4, 5.8, fill="white")
        add_text(slide, num, x + 0.45, 4.75, 1.1, 0.8, size=25, bold=True, color="blue", align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.55, 6.05, 4.3, 0.44, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.55, 6.9, 4.3, 1.2, size=9.6, color="muted", align=PP_ALIGN.CENTER)
    add_text(slide, "闭环价值：游客获得“随问随答、边走边讲”的服务，景区获得可分析、可维护、可持续优化的 AI 导览能力。", 2.2, 12.55, 29.2, 0.8, size=15, bold=True, align=PP_ALIGN.CENTER)


def function_matrix():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "功能完整性：覆盖赛题核心要求")
    rows = [
        ("多模态交互", "文本、浏览器语音、情绪标签、视觉理解接口", "已完成"),
        ("智能问答讲解", "DeepSeek 生成 + 本地景区知识库 RAG + 来源追踪", "已完成"),
        ("个性化推荐", "按游玩时长、兴趣偏好、景点热度推荐路线", "已完成"),
        ("数字人导览", "LiveTalking/Wav2Lip 驱动口型，WebRTC 前端播放", "已完成"),
        ("后台管理", "知识库、形象、反馈、行为数据分析、系统配置", "已完成"),
        ("数据看板", "游客画像、满意度、停留时长、景点热度统计", "已完成"),
        ("本地部署", "前端、后端、数字人服务均已在本机验证", "已完成"),
    ]
    add_card(slide, 1.45, 2.8, 30.9, 11.7, fill="white")
    y = 3.45
    widths = [7.0, 16.8, 4.0]
    headers = ["模块", "实现内容", "状态"]
    x0 = 2.0
    for i, htxt in enumerate(headers):
        add_text(slide, htxt, x0 + sum(widths[:i]), y, widths[i], 0.4, size=10.5, bold=True, color="muted")
    for r, (a, b, c) in enumerate(rows):
        yy = y + 0.9 + r * 1.34
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.95), Cm(yy - 0.18), Cm(29.8), Cm(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["line"]
        line.line.fill.background()
        add_text(slide, a, x0, yy, widths[0], 0.45, size=11.2, bold=True)
        add_text(slide, b, x0 + widths[0], yy, widths[1], 0.45, size=10.2, color="muted")
        add_pill(slide, c, x0 + widths[0] + widths[1] + 0.15, yy - 0.1, 2.2, "green")


def visitor():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "游客端体验：从问答到路线的完整旅程")
    add_image(slide, "lingling-guide-avatar.png", 2.0, 3.05, 7.0, 7.0)
    add_card(slide, 10.2, 3.0, 9.8, 8.8, "智能讲解", "游客提问“灵山大佛有什么特色？”后，系统检索本地知识库并调用大模型生成分点讲解，回答中保留来源信息，便于事实核验。", accent="blue")
    add_card(slide, 21.2, 3.0, 9.8, 8.8, "路线推荐", "输入游玩时长和兴趣偏好后，推荐“灵山大照壁 -> 祥符禅寺 -> 灵山大佛”等路线，并给出停留建议、总时长与推荐理由。", accent="teal")
    add_metric(slide, "本地问答实测延迟", "2.26s", 7.1, 13.1, 5.0, "green")
    add_metric(slide, "RAG 置信度示例", "0.90", 14.0, 13.1, 5.0, "blue")
    add_metric(slide, "核心接口 SLA 目标", "<5s", 20.9, 13.1, 5.0, "purple")


def digital_human():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "数字人导览：让 AI 回答变成可感知的讲解服务")
    add_image(slide, "scenic/photos/buddhist-culture-museum.jpg", 1.55, 3.0, 13.8, 8.2)
    add_card(slide, 16.3, 3.0, 15.2, 8.2, "技术实现", "数字人服务运行在 8010 端口，加载 Wav2Lip 模型进行口型驱动，使用 WebRTC 将视频流推送到前端；音色支持预设切换，并与游客端讲解内容联动。", accent="purple")
    add_bullets(slide, [
        "已验证页面：/index.html、/avatar.html",
        "已验证接口：/api/preset_voices、/api/admin/config、/api/avatar/tasks",
        "异常场景：无会话调用 /human 返回可控错误，不影响主系统运行",
    ], 17.0, 7.3, 13.8, 2.2, size=10.4)
    add_text(slide, "答辩展示建议：先用文本问答展示事实准确性，再切换数字人播报，强化“多模态 AI 导览助手”的观感记忆点。", 2.2, 13.15, 29.0, 0.7, size=14, bold=True, align=PP_ALIGN.CENTER)


def rag():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "本地知识库 RAG：提升准确性与可解释性")
    add_card(slide, 1.5, 3.0, 8.9, 8.6, "数据来源", "景区官方文档、结构化景点资料、行为分析数据和后台维护内容统一进入本地知识库。", accent="teal")
    add_card(slide, 12.3, 3.0, 8.9, 8.6, "检索增强", "先召回相关景点、路线和文档片段，再由大模型生成面向游客的自然语言回答。", accent="blue")
    add_card(slide, 23.1, 3.0, 8.9, 8.6, "事实约束", "回答返回置信度与来源，后台可持续补充知识，降低模型幻觉风险。", accent="green")
    add_metric(slide, "活跃知识文档", "50", 5.3, 13.0, 5.2, "teal")
    add_metric(slide, "示例置信度", "0.90", 14.35, 13.0, 5.2, "blue")
    add_metric(slide, "赛题准确率目标", ">=90%", 23.4, 13.0, 5.2, "green")


def route_map():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "路线与位置能力：把“知道”变成“会带路”")
    add_image(slide, "scenic/guide-map-clean.png", 1.55, 2.8, 15.1, 10.2)
    add_image(slide, "scenic/lingshan-handdrawn-map.png", 17.35, 2.8, 14.1, 10.2)
    add_card(slide, 1.8, 14.0, 29.3, 2.0, "推荐样例", "180 分钟历史文化偏好：灵山大照壁 -> 祥符禅寺 -> 灵山大佛。系统可按偏好、时长和景点属性输出路线理由，后续可接入实时定位与拥堵数据动态调整。", accent="blue")


def admin():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "管理后台：让内容、形象、反馈和数据可运营")
    cards = [
        ("知识库管理", "新增、编辑、审核景区讲解材料，持续提升 RAG 命中率。", "blue"),
        ("数字人配置", "管理形象、音色、欢迎语和服务状态，适配不同景区活动。", "purple"),
        ("反馈闭环", "收集游客评分、投诉建议和问题热度，反向优化知识库。", "amber"),
        ("数据看板", "分析客流、停留时长、满意度和热门景点，辅助运营决策。", "teal"),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(slide, 1.55 + i * 7.8, 3.1, 6.8, 7.2, t, b, accent=c)
    add_metric(slide, "行为明细记录", "140447", 4.2, 12.4, 5.8, "blue")
    add_metric(slide, "匹配景区记录", "777", 11.3, 12.4, 5.8, "teal")
    add_metric(slide, "平均满意度", "3.07", 18.4, 12.4, 5.8, "amber")
    add_metric(slide, "平均停留时长", "4.01h", 25.5, 12.4, 5.8, "green")


def architecture():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "系统架构：前后端、模型服务与数字人服务解耦")
    layers = [
        ("游客端 / 管理端", "Vue 3 + Vite，移动端入口、后台管理、数据可视化", "blue"),
        ("业务 API", "Python 后端，统一承载问答、路线、知识库、反馈和统计接口", "teal"),
        ("AI 能力层", "DeepSeek 文本生成、DashScope Qwen3-VL 多模态、RAG 检索增强", "purple"),
        ("数字人服务", "LiveTalking / Wav2Lip / WebRTC，负责口型驱动与流式播放", "amber"),
        ("数据层", "SQLite、本地文档、景区行为数据、前端静态资源", "green"),
    ]
    for i, (t, b, c) in enumerate(layers):
        y = 3.0 + i * 2.25
        add_card(slide, 3.0, y, 27.8, 1.55, t, b, accent=c)
        if i < len(layers) - 1:
            add_text(slide, "↓", 16.2, y + 1.46, 1.0, 0.45, size=16, bold=True, color="muted", align=PP_ALIGN.CENTER)
    add_text(slide, "工程优势：业务系统和数字人可独立启动、独立扩展，答辩演示时任一模块异常不会拖垮整体体验。", 2.3, 15.3, 29.0, 0.55, size=13.5, bold=True, align=PP_ALIGN.CENTER)


def innovation():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "技术创新点：不是单一聊天机器人，而是景区 AI 服务中台")
    items = [
        ("多模型协同", "文本模型负责生成，多模态模型理解图片/视觉输入，数字人模型负责可视化讲解。"),
        ("RAG 可控回答", "本地知识召回、置信度和来源回传，使景区内容可审计、可维护。"),
        ("数据驱动运营", "游客行为数据进入看板，支持热点发现、内容补强和服务排班。"),
        ("模块化部署", "前端、后端、数字人三服务解耦，适合景区从试点逐步扩展。"),
    ]
    for i, (t, b) in enumerate(items):
        add_card(slide, 2.0 + (i % 2) * 15.4, 3.2 + (i // 2) * 5.4, 13.8, 4.5, t, b, accent=["blue", "green", "teal", "purple"][i])
    add_text(slide, "可扩展方向：接入景区 IoT 客流、AR 导览、票务/餐饮/文创推荐，让导览助手从“讲解工具”升级为“全域服务入口”。", 2.1, 14.65, 29.6, 0.8, size=14, bold=True, align=PP_ALIGN.CENTER)


def verification():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "本地运行验证：关键服务与接口已打通")
    rows = [
        ("后端 API", "http://127.0.0.1:8000", "能力接口、问答、路线、数据看板均通过"),
        ("前端页面", "http://127.0.0.1:5173", "Vite 开发服务可访问，构建成功"),
        ("数字人服务", "http://127.0.0.1:8010", "页面、音色、配置、任务接口可访问"),
        ("移动端入口", "http://127.0.0.1:8000/mobile/", "移动端页面返回正常"),
    ]
    add_card(slide, 1.6, 3.0, 30.7, 8.2, fill="white")
    for i, (svc, url, result) in enumerate(rows):
        y = 3.6 + i * 1.65
        add_text(slide, svc, 2.2, y, 5.5, 0.42, size=11.5, bold=True)
        add_text(slide, url, 8.0, y, 8.0, 0.42, size=10.5, color="blue")
        add_text(slide, result, 17.0, y, 12.8, 0.42, size=10.2, color="muted")
        add_pill(slide, "通过", 29.0, y - 0.11, 1.8, "green")
    add_metric(slide, "后端单元测试", "55 OK", 4.7, 12.6, 5.8, "green")
    add_metric(slide, "前端单元测试", "4 OK", 11.8, 12.6, 5.8, "green")
    add_metric(slide, "前端构建", "PASS", 18.9, 12.6, 5.8, "blue")
    add_metric(slide, "数字人页面", "200", 26.0, 12.6, 5.8, "purple")


def tests():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "质量保障：用可复现测试支撑答辩可信度")
    add_card(slide, 1.55, 3.0, 9.4, 8.2, "自动化测试", "后端 55 个 unittest 用例通过；前端 Vitest 2 个测试文件、4 个用例通过；前端 typecheck 与 build 均成功。", accent="green")
    add_card(slide, 12.2, 3.0, 9.4, 8.2, "接口验证", "问答、路线推荐、能力接口、数据看板、移动端入口和数字人管理接口均已本机访问验证。", accent="blue")
    add_card(slide, 22.85, 3.0, 9.4, 8.2, "风险控制", "修复 Python 3.8 兼容、依赖版本、数据文件选择和 DLL 路径问题，降低评审现场运行风险。", accent="amber")
    add_text(slide, "可复现实测命令：后端 unittest、前端 typecheck/test/build、数字人页面和接口 HTTP 检查。", 2.2, 13.4, 29.4, 0.7, size=14, bold=True, align=PP_ALIGN.CENTER)


def business_value():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "落地价值：游客满意、景区增效、内容可持续")
    values = [
        ("游客侧", "随问随答、沉浸讲解、少走回头路，降低信息获取成本。"),
        ("景区侧", "减少重复咨询，提升高峰期服务承载，沉淀可分析的游客需求。"),
        ("运营侧", "通过满意度、停留时长和热点景点优化活动、路线和讲解内容。"),
        ("商业侧", "可扩展票务、餐饮、文创、会员服务推荐，形成二次转化入口。"),
    ]
    for i, (t, b) in enumerate(values):
        add_card(slide, 2.0 + (i % 2) * 15.4, 3.1 + (i // 2) * 4.9, 13.8, 4.0, t, b, accent=["blue", "teal", "green", "purple"][i])
    add_text(slide, "项目不只满足比赛功能项，也具备景区试点部署的现实路径：先做智能问答与后台，再接入数字人、定位和营销服务。", 2.2, 14.2, 29.4, 0.8, size=14.5, bold=True, align=PP_ALIGN.CENTER)


def demo():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "2 分钟演示脚本：突出高分记忆点")
    steps = [
        ("0-20s", "打开游客端", "展示灵灵导游入口、语音/文本交互和景区视觉资产。"),
        ("20-50s", "问答讲解", "提问灵山大佛特色，展示 RAG 回答、置信度和来源。"),
        ("50-75s", "路线推荐", "设置 180 分钟与历史文化偏好，展示路线和理由。"),
        ("75-100s", "数字人播报", "切到 8010 数字人页面，播放同一讲解内容。"),
        ("100-120s", "后台数据", "展示知识库、反馈和行为看板，强调运营闭环。"),
    ]
    for i, (time, title, body) in enumerate(steps):
        add_card(slide, 2.0, 3.0 + i * 2.25, 29.2, 1.55, f"{time}  {title}", body, accent=["blue", "blue", "teal", "purple", "green"][i])
    add_text(slide, "讲解口径：每一步都回扣赛题评分项，避免只演示页面而没有解释技术和落地价值。", 2.2, 15.0, 29.2, 0.55, size=13.5, bold=True, align=PP_ALIGN.CENTER)


def delivery():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "交付清单与后续优化")
    add_card(slide, 1.6, 3.0, 14.3, 9.3, "本次可提交材料", "源代码与可执行脚本、部署运行说明、数据库和知识库资料、设计说明文档、高分版 PPT、2 分钟演示视频脚本。", accent="blue")
    add_card(slide, 17.3, 3.0, 14.3, 9.3, "下一步加分优化", "接入实时定位和地图导航；补充更多景区官方内容；扩展多语言讲解；增加票务/餐饮/文创推荐；上线日志监控和服务健康检查。", accent="green")
    add_text(slide, "总结：本项目以完整功能闭环、可解释 RAG、多模态数字人和可运行工程交付，形成符合 A5 赛题要求的高完成度方案。", 2.2, 14.3, 29.3, 0.9, size=15, bold=True, align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]

for fn in [
    cover,
    scoring,
    pain_points,
    loop,
    function_matrix,
    visitor,
    digital_human,
    rag,
    route_map,
    admin,
    architecture,
    innovation,
    verification,
    tests,
    business_value,
    demo,
    delivery,
]:
    fn()

DOCS.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
print(len(prs.slides))
