# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[2]
SOFTWAREBAR = ROOT / "softwarebar"
ASSETS = SOFTWAREBAR / "frontend-vue" / "public" / "assets"
DOCS = SOFTWAREBAR / "docs"
OUT = DOCS / "产品方案介绍PPT_答辩优化版.pptx"

SLIDE_W, SLIDE_H = Cm(33.867), Cm(19.05)

COLORS = {
    "ink": RGBColor(20, 30, 43),
    "muted": RGBColor(91, 104, 124),
    "soft": RGBColor(246, 248, 251),
    "line": RGBColor(219, 228, 238),
    "white": RGBColor(255, 255, 255),
    "dark": RGBColor(12, 19, 31),
    "blue": RGBColor(38, 99, 235),
    "teal": RGBColor(13, 148, 136),
    "green": RGBColor(22, 163, 74),
    "amber": RGBColor(217, 119, 6),
    "red": RGBColor(220, 38, 38),
    "slate": RGBColor(51, 65, 85),
}


def rgb(name):
    return COLORS[name] if isinstance(name, str) else name


def set_bg(slide, color="soft"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=18, bold=False, color="ink", align=None, valign=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_multi(slide, lines, x, y, w, h, size=17, color="ink", bullet=False, gap=8):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
        if bullet:
            p.level = 0
            p.text = "• " + line
    return box


def title(slide, main, sub=None, dark=False):
    color = "white" if dark else "ink"
    sub_color = RGBColor(214, 225, 240) if dark else COLORS["muted"]
    add_text(slide, main, 1.45, 0.72, 26.5, 0.9, size=27, bold=True, color=color)
    if sub:
        add_text(slide, sub, 1.5, 1.68, 28.2, 0.52, size=13, color=sub_color)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.45), Cm(2.35), Cm(30.9), Cm(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(80, 96, 120) if dark else COLORS["line"]
    line.line.fill.background()


def card(slide, x, y, w, h, fill="white", line="line", radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.8)
    shp.adjustments[0] = radius
    return shp


def band(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    return shp


def metric(slide, value, label, x, y, w, accent="blue", note=None):
    card(slide, x, y, w, 2.85)
    add_text(slide, value, x + 0.2, y + 0.42, w - 0.4, 0.8, size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.35, y + 1.38, w - 0.7, 0.48, size=13, bold=True, color="ink", align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + 0.35, y + 2.05, w - 0.7, 0.35, size=10.5, color="muted", align=PP_ALIGN.CENTER)


def pill(slide, text, x, y, w, color="blue"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.62))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    shp.adjustments[0] = 0.42
    add_text(slide, text, x, y + 0.11, w, 0.34, size=10, bold=True, color="white", align=PP_ALIGN.CENTER)


def image(slide, rel, x, y, w, h):
    path = ASSETS / rel
    if path.exists():
        slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    else:
        card(slide, x, y, w, h, fill="white")
        add_text(slide, f"缺少图片\n{rel}", x + 0.5, y + h / 2 - 0.5, w - 1.0, 1.0, size=14, color="muted", align=PP_ALIGN.CENTER)


def footer(slide, page):
    add_text(slide, f"{page:02d}", 31.5, 17.55, 0.8, 0.35, size=10, color="muted", align=PP_ALIGN.RIGHT)


def icon_circle(slide, label, x, y, color="blue"):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(1.45), Cm(1.45))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    add_text(slide, label, x, y + 0.34, 1.45, 0.45, size=16, bold=True, color="white", align=PP_ALIGN.CENTER)


def cover():
    s = prs.slides.add_slide(blank)
    set_bg(s, "dark")
    image(s, "scenic/landing-hero.png", 14.7, 0, 19.2, 19.05)
    overlay = band(s, 0, 0, 19.2, 19.05, "dark")
    overlay.fill.transparency = 6
    add_text(s, "A5 多模态 AI\n智能景区导览助手", 1.7, 2.5, 14.2, 2.4, size=35, bold=True, color="white")
    add_text(s, "灵灵导游：数字人讲解 + 本地知识库 RAG + 个性化路线 + 运营后台", 1.78, 5.45, 15.3, 0.6, size=15, color=RGBColor(224, 235, 250))
    for i, (t, c) in enumerate([("可运行", "green"), ("可验证", "blue"), ("可落地", "teal")]):
        pill(s, t, 1.78 + i * 3.0, 6.55, 2.35, c)
    add_text(s, "参赛作品答辩版", 1.78, 16.95, 5.2, 0.42, size=12, color=RGBColor(210, 220, 234))
    footer(s, 1)


def score_map():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "评分导向：每一页都回答“为什么能拿分”", "把赛题要求翻译成可演示、可验证、可落地的四类证据")
    items = [
        ("40%", "功能完整性", "游客端、数字人、路线推荐、后台管理形成闭环", "blue"),
        ("30%", "技术创新", "DeepSeek + Qwen-VL + RAG + Wav2Lip 多模型协同", "teal"),
        ("20%", "行业价值", "景区服务提效，游客体验提升，运营数据沉淀", "green"),
        ("10%", "文档展示", "部署、测试、PPT、视频脚本均可直接交付", "amber"),
    ]
    for i, (pct, name, desc, color) in enumerate(items):
        x = 1.65 + i * 7.85
        card(s, x, 4.1, 6.95, 7.5)
        band(s, x, 4.1, 6.95, 0.28, color)
        add_text(s, pct, x + 0.35, 4.95, 6.25, 0.95, size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, name, x + 0.45, 6.35, 6.05, 0.48, size=17, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, desc, x + 0.65, 7.42, 5.65, 1.35, size=13.5, color="muted", align=PP_ALIGN.CENTER)
        add_text(s, "本项目已覆盖", x + 1.15, 9.55, 4.65, 0.45, size=12.5, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, "答辩策略：先证明“功能完整”，再证明“技术可信”，最后落到“景区能用”。", 2.15, 14.15, 29.5, 0.72, size=20, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 2)


def positioning():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "产品定位：不是聊天窗口，而是景区 AI 服务入口")
    image(s, "lingling-avatar-v2.png", 1.55, 3.1, 8.7, 10.9)
    add_text(s, "灵灵导游", 2.05, 14.1, 7.5, 0.55, size=20, bold=True, color="blue", align=PP_ALIGN.CENTER)
    columns = [
        ("游客侧", ["问景点", "听讲解", "要路线"], "blue"),
        ("AI 侧", ["RAG 检索", "大模型生成", "多模态理解"], "teal"),
        ("景区侧", ["维护知识", "看反馈", "看数据"], "green"),
    ]
    for i, (head, lines, color) in enumerate(columns):
        x = 11.3 + i * 6.9
        card(s, x, 4.0, 5.95, 8.8)
        add_text(s, head, x + 0.35, 4.65, 5.25, 0.52, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, item in enumerate(lines):
            add_text(s, item, x + 0.78, 6.1 + j * 1.55, 4.4, 0.5, size=18, bold=True, align=PP_ALIGN.CENTER)
        pill(s, "形成闭环", x + 1.72, 11.4, 2.5, color)
    add_text(s, "一个入口同时服务游客体验、AI 生成和景区运营，避免作品停留在单点功能。", 11.4, 14.45, 20.0, 0.8, size=19, bold=True)
    footer(s, 3)


def product_loop():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "业务闭环：问、讲、游、管全部打通")
    steps = [
        ("问", "文本 / 语音 / 情绪", "游客提出问题"),
        ("查", "本地知识库 RAG", "召回景区事实"),
        ("讲", "数字人 + TTS", "沉浸式播报"),
        ("游", "偏好路线推荐", "动态规划行程"),
        ("管", "反馈与数据看板", "运营持续优化"),
    ]
    for i, (big, name, desc) in enumerate(steps):
        x = 1.55 + i * 6.35
        icon_circle(s, big, x + 1.95, 4.1, ["blue", "teal", "amber", "green", "slate"][i])
        add_text(s, name, x + 0.2, 6.25, 5.1, 0.5, size=17, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, desc, x + 0.35, 7.1, 4.8, 0.45, size=13.5, color="muted", align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(s, "→", x + 5.3, 4.5, 1.0, 0.6, size=28, bold=True, color="muted", align=PP_ALIGN.CENTER)
    card(s, 3.0, 11.0, 27.8, 2.85, fill="white")
    add_text(s, "核心价值", 4.0, 11.42, 4.0, 0.5, size=17, bold=True, color="blue")
    add_text(s, "游客得到“随问随答、边走边讲”的服务；景区得到可维护、可分析、可持续优化的 AI 导览能力。", 8.1, 11.4, 21.2, 0.95, size=19, bold=True)
    footer(s, 4)


def visitor_experience():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "游客端体验：3 步完成一次智能导览")
    image(s, "scenic/photos/lingshan-grand-buddha.jpg", 1.55, 3.0, 12.6, 9.45)
    panels = [
        ("1", "提问", "“灵山大佛有什么特色？”"),
        ("2", "回答", "RAG 引用本地资料，生成分点讲解"),
        ("3", "行动", "按 180 分钟与历史文化偏好推荐路线"),
    ]
    for i, (num, head, body) in enumerate(panels):
        y = 3.05 + i * 3.1
        card(s, 15.3, y, 15.9, 2.4)
        icon_circle(s, num, 15.85, y + 0.48, ["blue", "teal", "green"][i])
        add_text(s, head, 17.65, y + 0.42, 4.5, 0.5, size=18, bold=True)
        add_text(s, body, 17.65, y + 1.18, 12.8, 0.55, size=15.5, color="muted")
    metric(s, "2.26s", "问答实测延迟", 15.3, 12.9, 4.9, "green")
    metric(s, "0.90", "RAG 置信度", 21.0, 12.9, 4.9, "blue")
    metric(s, "<5s", "赛题稳定性目标", 26.7, 12.9, 4.9, "teal")
    footer(s, 5)


def feature_coverage():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "功能完整性：赛题核心功能逐项落地")
    rows = [
        ("多模态交互", "文本、语音、情绪、视觉理解接口", "已实现"),
        ("智能讲解问答", "DeepSeek + 本地知识库 RAG + 来源追踪", "已实现"),
        ("个性化路线", "按时长、偏好、景点属性推荐路径", "已实现"),
        ("数字人导览", "Wav2Lip 口型驱动 + WebRTC 播放", "已实现"),
        ("后台管理", "知识库、形象、反馈、配置、数据看板", "已实现"),
        ("本地部署", "前端、后端、数字人三服务本机跑通", "已验证"),
    ]
    card(s, 1.65, 3.05, 30.6, 11.7)
    y = 3.65
    for i, (m, d, status) in enumerate(rows):
        yy = y + i * 1.72
        if i:
            band(s, 2.2, yy - 0.26, 29.3, 0.025, "line")
        add_text(s, m, 2.45, yy, 6.2, 0.55, size=16.5, bold=True)
        add_text(s, d, 9.2, yy, 16.6, 0.55, size=15, color="muted")
        pill(s, status, 27.0, yy - 0.05, 2.55, "green")
    add_text(s, "展示重点：评委能直接看到“必做项没有缺口”。", 2.15, 15.45, 29.5, 0.55, size=18.5, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 6)


def rag_accuracy():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "RAG 让回答更准确：先查本地资料，再由模型组织语言")
    stages = [
        ("景区资料", "官方文档 / 景点表 / 后台维护"),
        ("向量与关键词召回", "定位相关景点、路线和讲解片段"),
        ("大模型生成", "组织成游客能听懂的讲解"),
        ("来源与置信度", "给出可核验依据，降低幻觉风险"),
    ]
    for i, (head, body) in enumerate(stages):
        x = 1.65 + i * 7.75
        card(s, x, 4.15, 6.65, 5.7)
        icon_circle(s, str(i + 1), x + 2.6, 4.75, ["blue", "teal", "green", "amber"][i])
        add_text(s, head, x + 0.55, 6.65, 5.55, 0.5, size=16.2, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, body, x + 0.65, 7.62, 5.35, 0.95, size=13.8, color="muted", align=PP_ALIGN.CENTER)
    metric(s, "50", "活跃知识文档", 6.2, 12.2, 5.5, "teal")
    metric(s, "0.90", "示例回答置信度", 14.2, 12.2, 5.5, "blue")
    metric(s, "≥90%", "赛题准确率目标", 22.2, 12.2, 5.5, "green")
    footer(s, 7)


def digital_human():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "数字人导览：把文字回答变成可感知的讲解服务")
    image(s, "lingling-guide-avatar.png", 1.8, 3.1, 9.2, 9.2)
    card(s, 12.3, 3.2, 18.9, 8.9)
    add_text(s, "技术链路", 13.1, 3.85, 5.0, 0.5, size=18, bold=True, color="blue")
    add_multi(s, [
        "Wav2Lip 加载口型模型，驱动头像说话",
        "EdgeTTS / 预设音色输出语音",
        "WebRTC 推送音视频流到浏览器",
        "与游客端问答内容联动，可用于现场演示",
    ], 13.15, 4.75, 16.6, 4.5, size=16.2, color="ink", bullet=True, gap=10)
    metric(s, "8010", "数字人服务端口", 12.3, 13.0, 5.6, "blue")
    metric(s, "200", "页面与接口状态", 19.1, 13.0, 5.6, "green")
    metric(s, "1", "演示主会话", 25.9, 13.0, 5.6, "teal")
    footer(s, 8)


def route_map():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "路线推荐：从“会回答”升级到“会带路”")
    image(s, "scenic/guide-map-clean.png", 1.6, 3.0, 15.0, 10.8)
    card(s, 17.5, 3.05, 13.8, 10.75)
    add_text(s, "推荐样例", 18.25, 3.8, 5.0, 0.5, size=18, bold=True, color="teal")
    add_text(s, "180 分钟 · 历史文化偏好", 18.25, 4.8, 9.0, 0.55, size=18, bold=True)
    route = ["灵山大照壁", "祥符禅寺", "灵山大佛"]
    for i, name in enumerate(route):
        y = 6.25 + i * 1.65
        icon_circle(s, str(i + 1), 18.35, y - 0.32, ["blue", "teal", "green"][i])
        add_text(s, name, 20.15, y, 8.8, 0.48, size=17, bold=True)
    add_text(s, "路线依据：游玩时长、兴趣偏好、景点属性与热度。后续可接入实时定位和拥堵数据，实现动态导览。", 18.25, 11.75, 11.4, 1.0, size=14.8, color="muted")
    footer(s, 9)


def admin_data():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "后台运营：景区能维护内容，也能看见游客行为")
    image(s, "scenic/photos/buddhist-culture-museum.jpg", 1.55, 3.05, 10.6, 9.1)
    funcs = [("知识库", "持续补充讲解内容"), ("形象配置", "管理数字人与音色"), ("反馈闭环", "沉淀问题与评分"), ("行为看板", "分析客流与偏好")]
    for i, (head, body) in enumerate(funcs):
        x = 13.2 + (i % 2) * 8.75
        y = 3.1 + (i // 2) * 3.0
        card(s, x, y, 7.65, 2.35)
        add_text(s, head, x + 0.45, y + 0.45, 6.75, 0.45, size=16.5, bold=True, color=["blue", "teal", "green", "amber"][i])
        add_text(s, body, x + 0.45, y + 1.25, 6.55, 0.45, size=13.5, color="muted")
    metric(s, "140447", "行为明细记录", 13.2, 10.15, 5.3, "blue")
    metric(s, "777", "匹配景区记录", 19.55, 10.15, 5.3, "teal")
    metric(s, "4.01h", "平均停留时长", 25.9, 10.15, 5.3, "green")
    add_text(s, "运营价值：游客互动不再只是一次问答，而会进入内容优化和服务决策闭环。", 2.2, 15.0, 29.5, 0.55, size=18.5, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 10)


def architecture():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "系统架构：三服务解耦，便于本地部署和现场演示")
    layers = [
        ("前端体验层", "Vue 3 / Vite / 游客端 / 管理端", "blue"),
        ("业务 API 层", "Python 后端 / 问答 / 路线 / 知识库 / 数据看板", "teal"),
        ("AI 能力层", "DeepSeek / Qwen3-VL / RAG 检索增强", "green"),
        ("数字人服务", "LiveTalking / Wav2Lip / EdgeTTS / WebRTC", "amber"),
        ("数据资源层", "SQLite / 景区文档 / 行为数据 / 静态资源", "slate"),
    ]
    for i, (head, body, color) in enumerate(layers):
        y = 3.0 + i * 2.35
        card(s, 4.1, y, 25.6, 1.55)
        band(s, 4.1, y, 0.32, 1.55, color)
        add_text(s, head, 5.1, y + 0.36, 6.0, 0.45, size=16.5, bold=True, color=color)
        add_text(s, body, 11.4, y + 0.37, 16.7, 0.45, size=15.2, color="muted")
    footer(s, 11)


def innovation():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "技术创新：多模型协同 + 可控知识 + 可运营数据")
    items = [
        ("多模型协同", "文本生成、视觉理解、数字人口型驱动分别由最合适的模型承担"),
        ("RAG 可控回答", "本地资料优先，输出来源与置信度，避免泛泛聊天"),
        ("服务解耦部署", "前端、后端、数字人独立运行，演示与扩展风险更低"),
        ("数据反哺运营", "反馈和行为数据进入看板，持续优化讲解与路线"),
    ]
    for i, (head, body) in enumerate(items):
        x = 2.0 + (i % 2) * 15.25
        y = 3.4 + (i // 2) * 5.0
        card(s, x, y, 13.5, 3.95)
        add_text(s, head, x + 0.55, y + 0.55, 11.9, 0.55, size=18, bold=True, color=["blue", "teal", "green", "amber"][i])
        add_text(s, body, x + 0.6, y + 1.55, 12.15, 1.05, size=15, color="muted")
    add_text(s, "这使作品从“AI 问答 Demo”提升为“景区智能导览解决方案”。", 2.4, 14.3, 29.0, 0.65, size=20, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 12)


def verification():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "本地验证：不是概念稿，已经能在电脑上跑起来")
    checks = [
        ("后端 API", "8000", "能力接口、问答、路线、数据看板通过"),
        ("前端页面", "5173", "typecheck / unit test / build 通过"),
        ("移动端入口", "8000/mobile", "移动页面访问正常"),
        ("数字人服务", "8010", "页面、音色、会话接口通过"),
    ]
    for i, (svc, port, desc) in enumerate(checks):
        y = 3.2 + i * 2.15
        card(s, 2.0, y, 29.6, 1.55)
        add_text(s, svc, 2.8, y + 0.38, 5.2, 0.45, size=16.5, bold=True)
        add_text(s, port, 9.0, y + 0.35, 4.4, 0.45, size=17, bold=True, color="blue")
        add_text(s, desc, 14.2, y + 0.4, 12.5, 0.42, size=14.5, color="muted")
        pill(s, "PASS", 28.0, y + 0.45, 2.3, "green")
    metric(s, "55 OK", "后端测试", 6.0, 12.9, 5.4, "green")
    metric(s, "4 OK", "前端测试", 14.2, 12.9, 5.4, "green")
    metric(s, "200", "关键页面状态", 22.4, 12.9, 5.4, "blue")
    footer(s, 13)


def demo_script():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "2 分钟演示安排：让评委看到完整闭环")
    timeline = [
        ("0-20s", "打开游客端", "展示灵灵导游入口"),
        ("20-50s", "问答讲解", "提问灵山大佛，展示 RAG 回答"),
        ("50-75s", "路线推荐", "输入时长与偏好，得到路线"),
        ("75-100s", "数字人播报", "连接 8010，展示口型驱动"),
        ("100-120s", "后台看板", "展示知识库、反馈、行为数据"),
    ]
    for i, (time, head, body) in enumerate(timeline):
        x = 1.6 + i * 6.25
        card(s, x, 4.0, 5.4, 8.0)
        add_text(s, time, x + 0.45, 4.6, 4.5, 0.48, size=15, bold=True, color=["blue", "teal", "green", "amber", "slate"][i], align=PP_ALIGN.CENTER)
        add_text(s, head, x + 0.35, 6.0, 4.7, 0.58, size=17, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, body, x + 0.55, 7.15, 4.3, 1.25, size=13.6, color="muted", align=PP_ALIGN.CENTER)
    add_text(s, "讲解口径：每一步都回扣评分项，避免只演示页面而没有解释技术和价值。", 2.2, 14.2, 29.2, 0.65, size=18.5, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 14)


def closing():
    s = prs.slides.add_slide(blank)
    set_bg(s, "dark")
    image(s, "scenic/photos/bodhi-avenue.png", 17.2, 0, 16.7, 19.05)
    band(s, 0, 0, 22.2, 19.05, "dark")
    add_text(s, "总结：高完成度、可运行、可落地", 1.75, 2.25, 18.2, 0.95, size=31, bold=True, color="white")
    add_multi(s, [
        "功能闭环完整：游客端、数字人、路线、后台都能演示",
        "技术路径可信：RAG、多模型协同、WebRTC 数字人",
        "行业价值清晰：提升游客体验，沉淀景区运营数据",
        "交付材料齐全：代码、部署说明、PPT、演示脚本可提交",
    ], 2.0, 4.25, 16.8, 5.6, size=17.5, color=RGBColor(226, 236, 248), bullet=True, gap=12)
    for i, (v, l, c) in enumerate([("40%", "功能", "blue"), ("30%", "技术", "teal"), ("20%", "落地", "green"), ("10%", "展示", "amber")]):
        metric(s, v, l, 2.0 + i * 4.0, 12.3, 3.25, c)
    footer(s, 15)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

for builder in [
    cover,
    score_map,
    positioning,
    product_loop,
    visitor_experience,
    feature_coverage,
    rag_accuracy,
    digital_human,
    route_map,
    admin_data,
    architecture,
    innovation,
    verification,
    demo_script,
    closing,
]:
    builder()

DOCS.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
print(len(prs.slides))
