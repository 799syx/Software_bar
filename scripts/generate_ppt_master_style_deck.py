# -*- coding: utf-8 -*-
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[2]
ASSETS = ROOT / "softwarebar" / "frontend-vue" / "public" / "assets"
DOCS = ROOT / "softwarebar" / "docs"
OUT = DOCS / "灵灵导游_A5多模态AI智能景区导览助手_正式展示亮色版_含数据来源.pptx"

SLIDE_W, SLIDE_H = Cm(33.867), Cm(19.05)
W_CM, H_CM = 33.867, 19.05
TOTAL_SLIDES = 16

COLORS = {
    "bg": RGBColor(245, 249, 252),
    "panel": RGBColor(255, 255, 255),
    "glass": RGBColor(248, 252, 255),
    "white": RGBColor(255, 255, 255),
    "text": RGBColor(28, 43, 63),
    "muted": RGBColor(91, 110, 132),
    "ink": RGBColor(18, 31, 48),
    "blue": RGBColor(56, 189, 248),
    "teal": RGBColor(45, 212, 191),
    "green": RGBColor(34, 197, 94),
    "amber": RGBColor(245, 158, 11),
    "violet": RGBColor(167, 139, 250),
    "line": RGBColor(205, 218, 232),
    "red": RGBColor(248, 113, 113),
}


def rgb(name):
    return COLORS[name] if isinstance(name, str) else name


def blank_slide(prs, dark=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg" if dark else "white")
    return slide


def set_text_frame(tf):
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.08)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)


def add_text(slide, text, x, y, w, h, size=18, bold=False, color="text", align=None, valign=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    set_text_frame(tf)
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_lines(slide, lines, x, y, w, h, size=17, color="text", gap=8, bullet=False):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    set_text_frame(tf)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if bullet else "") + line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
    return box


def rect(slide, x, y, w, h, fill="panel", line=None, transparency=0, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def glow(slide, x, y, w, h, color="blue", transparency=78):
    # Avoid decorative glow orbs: they render inconsistently across PowerPoint/LibreOffice.
    return None


def cover_picture(slide, rel, x, y, w, h):
    path = ASSETS / rel
    if not path.exists():
        rect(slide, x, y, w, h, "panel", "line", 0, True)
        add_text(slide, f"缺少图片：{rel}", x + 0.4, y + h / 2 - 0.3, w - 0.8, 0.6, 15, color="muted", align=PP_ALIGN.CENTER)
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = w / h
    if img_ratio >= box_ratio:
        pic = slide.shapes.add_picture(str(path), Cm(x), Cm(y), height=Cm(h))
        excess = (pic.width - Cm(w)) / pic.width
        pic.crop_left = excess / 2
        pic.crop_right = excess / 2
        pic.width = Cm(w)
    else:
        pic = slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
        excess = (pic.height - Cm(h)) / pic.height
        pic.crop_top = excess / 2
        pic.crop_bottom = excess / 2
        pic.height = Cm(h)
    pic.left = Cm(x)
    pic.top = Cm(y)
    return pic


def header(slide, title, kicker=None, page=1, dark=True):
    add_text(slide, f"{page:02d}", 1.25, 0.65, 1.0, 0.45, 11, bold=True, color="blue")
    add_text(slide, title, 2.35, 0.55, 24.5, 0.78, 24, bold=True, color="text" if dark else "ink")
    if kicker:
        add_text(slide, kicker, 2.38, 1.42, 24.5, 0.44, 11.5, color="muted")
    rect(slide, 1.25, 2.13, 31.35, 0.03, "line", transparency=35)


def footer(slide, page):
    add_text(slide, "灵灵导游 · A5 多模态 AI 智能景区导览助手", 1.25, 17.92, 11, 0.36, 9.5, color="muted")
    add_text(slide, f"{page:02d}/{TOTAL_SLIDES}", 30.5, 17.92, 2.1, 0.36, 9.5, color="muted", align=PP_ALIGN.RIGHT)


def source_note(slide, text, x=1.25, y=16.95, w=31.3):
    add_text(slide, f"来源：{text}", x, y, w, 0.4, 8.8, color="muted")


def pill(slide, text, x, y, w, color="blue", light=False):
    fill = "white" if light else color
    shape = rect(slide, x, y, w, 0.62, fill, None, 0 if not light else 15, True)
    add_text(slide, text, x, y + 0.12, w, 0.34, 10.5, bold=True, color="ink" if light else "bg", align=PP_ALIGN.CENTER)
    return shape


def glass_panel(slide, x, y, w, h, alpha=23):
    rect(slide, x, y, w, h, "panel", "line", 0, True)


def big_metric(slide, value, label, x, y, w, color="blue", note=None):
    glass_panel(slide, x, y, w, 3.15, 18)
    add_text(slide, value, x + 0.25, y + 0.38, w - 0.5, 0.9, 30, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.35, y + 1.45, w - 0.7, 0.45, 13.5, bold=True, color="text", align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + 0.35, y + 2.25, w - 0.7, 0.35, 10.5, color="muted", align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color="line"):
    line = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2.2)
    return line


def circle_label(slide, label, x, y, color="blue", size=1.65):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(size), Cm(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    add_text(slide, label, x, y + size * 0.28, size, size * 0.36, 16, bold=True, color="bg", align=PP_ALIGN.CENTER)


def make_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 01 Cover
    s = blank_slide(prs)
    cover_picture(s, "scenic/landing-hero.png", 0, 0, W_CM, H_CM)
    rect(s, 0, 0, 19.4, H_CM, "bg", transparency=8)
    glow(s, -2.2, -1.8, 12, 12, "blue", 72)
    add_text(s, "A5 多模态 AI\n智能景区导览助手", 1.65, 2.55, 15.4, 2.8, 38, bold=True, color="ink")
    add_text(s, "灵灵导游：数字人讲解 × 本地知识库 RAG × 个性化路线 × 运营后台", 1.75, 5.92, 15.8, 0.6, 15.5, color="muted")
    for i, item in enumerate(["可运行", "可验证", "可落地"]):
        pill(s, item, 1.78 + i * 3.05, 6.95, 2.42, ["green", "blue", "teal"][i])
    add_text(s, "参赛作品展示版", 1.75, 16.8, 7.2, 0.42, 12.5, color="muted")

    # 02 Scoring map
    s = blank_slide(prs)
    header(s, "赛题要求：四类能力完整覆盖", "从功能、技术、行业与文档四个维度呈现项目完成度", 2)
    scores = [("40%", "功能完整性", "问答、讲解、路线、数字人、后台均可演示", "blue"),
              ("30%", "技术创新", "DeepSeek + Qwen-VL + RAG + Wav2Lip 多模型协同", "teal"),
              ("20%", "行业适用性", "游客体验提效，景区运营数据沉淀", "green"),
              ("10%", "文档展示", "源码、部署、测试、PPT、演示脚本闭环", "amber")]
    for i, (v, title, body, c) in enumerate(scores):
        x = 1.55 + i * 7.9
        glass_panel(s, x, 4.2, 6.9, 8.0, 18)
        rect(s, x, 4.2, 6.9, 0.24, c)
        add_text(s, v, x + 0.35, 5.05, 6.2, 0.95, 35, True, c, PP_ALIGN.CENTER)
        add_text(s, title, x + 0.45, 6.55, 6.0, 0.48, 17, True, "ink", PP_ALIGN.CENTER)
        add_text(s, body, x + 0.62, 7.65, 5.65, 1.25, 13.8, color="muted", align=PP_ALIGN.CENTER)
        pill(s, "已覆盖", x + 2.0, 10.45, 2.85, c)
    add_text(s, "展示内容覆盖完整能力、技术链路与景区落地价值。", 3.0, 14.25, 27.8, 0.7, 22, True, "ink", PP_ALIGN.CENTER)
    source_note(s, "赛题说明与评分结构：功能完整性 40%、技术与创新 30%、行业适用性 20%、文档与展示 10%。")
    footer(s, 2)

    # 03 Positioning
    s = blank_slide(prs)
    header(s, "项目定位：景区 AI 导览服务平台", "统一连接游客体验、AI 能力和景区运营", 3)
    cover_picture(s, "lingling-avatar-v2.png", 1.6, 3.15, 8.2, 10.3)
    add_text(s, "灵灵导游", 2.0, 13.85, 7.4, 0.55, 22, True, "blue", PP_ALIGN.CENTER)
    cols = [("游客侧", ["问景点", "听讲解", "要路线"], "blue"),
            ("AI 侧", ["RAG 检索", "多模型生成", "口型驱动"], "teal"),
            ("景区侧", ["维护知识", "看反馈", "看数据"], "green")]
    for i, (head, items, c) in enumerate(cols):
        x = 11.2 + i * 6.9
        glass_panel(s, x, 4.15, 5.95, 8.35, 21)
        add_text(s, head, x + 0.35, 4.9, 5.25, 0.55, 18, True, c, PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text(s, item, x + 0.65, 6.35 + j * 1.48, 4.7, 0.52, 19, True, "ink", PP_ALIGN.CENTER)
        pill(s, "形成闭环", x + 1.6, 11.05, 2.75, c)
    add_text(s, "游客咨询、知识检索、数字人讲解、路线推荐与后台运营形成完整服务链。", 11.45, 14.15, 20.0, 0.9, 21, True, "ink")
    source_note(s, "项目功能设计文档与本地系统页面：游客端、AI 能力层、景区后台三端功能。")
    footer(s, 3)

    # 04 Journey loop
    s = blank_slide(prs)
    header(s, "游客旅程闭环：问、查、讲、游、管", "把游客的一次咨询转化为完整导览服务", 4)
    steps = [("问", "文本/语音/情绪", "游客提出需求", "blue"),
             ("查", "本地知识库 RAG", "召回景区事实", "teal"),
             ("讲", "数字人播报", "口型驱动讲解", "amber"),
             ("游", "偏好路线推荐", "规划停留顺序", "green"),
             ("管", "反馈与看板", "运营持续优化", "violet")]
    for i, (n, t, b, c) in enumerate(steps):
        x = 1.45 + i * 6.35
        circle_label(s, n, x + 1.75, 4.1, c, 1.9)
        add_text(s, t, x + 0.1, 6.55, 5.45, 0.45, 17, True, "ink", PP_ALIGN.CENTER)
        add_text(s, b, x + 0.2, 7.35, 5.2, 0.42, 13.5, color="muted", align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(s, "→", x + 5.45, 4.42, 0.9, 0.55, 27, True, "muted", PP_ALIGN.CENTER)
    glass_panel(s, 3.0, 11.25, 27.85, 2.7, 20)
    add_text(s, "核心价值", 4.0, 11.74, 3.6, 0.5, 16, True, "blue")
    add_text(s, "游客获得随问随答、边走边讲的体验；景区获得可维护、可分析、可持续优化的 AI 导览能力。", 8.0, 11.68, 21.1, 0.72, 20, True, "ink")
    footer(s, 4)

    # 05 Visitor proof
    s = blank_slide(prs)
    header(s, "游客端体验：三步完成一次智能导览", "以真实问题和路线偏好呈现完整交互流程", 5)
    cover_picture(s, "scenic/photos/lingshan-grand-buddha.jpg", 1.35, 3.0, 12.7, 9.55)
    for i, (num, head, body, c) in enumerate([
        ("1", "提出问题", "“灵山大佛有什么特色？”", "blue"),
        ("2", "RAG 回答", "返回置信度、来源和分点讲解", "teal"),
        ("3", "生成路线", "180 分钟 + 历史文化偏好", "green")]):
        y = 3.1 + i * 2.95
        glass_panel(s, 15.2, y, 16.0, 2.22, 16)
        circle_label(s, num, 15.82, y + 0.43, c, 1.22)
        add_text(s, head, 17.45, y + 0.38, 5.8, 0.48, 17, True, "ink")
        add_text(s, body, 17.45, y + 1.15, 12.4, 0.45, 14.5, color="muted")
    big_metric(s, "2.26s", "问答实测延迟", 15.2, 12.75, 5.0, "green")
    big_metric(s, "0.90", "RAG 置信度", 20.7, 12.75, 5.0, "blue")
    big_metric(s, "<5s", "赛题响应目标", 26.2, 12.75, 5.0, "teal")
    source_note(s, "本地接口验证记录：示例问题“灵山大佛有什么特色？”返回置信度 0.90，响应延迟约 2.26 秒；赛题目标响应延迟不高于 5 秒。")
    footer(s, 5)

    # 06 Function matrix
    s = blank_slide(prs)
    header(s, "功能覆盖：赛题必做项逐项落地", "每个能力都有对应页面、接口或后台模块支撑", 6)
    rows = [("文字/语音/情绪交互", "游客端输入 + 浏览器语音能力 + 情绪/偏好标签", "已实现"),
            ("景点问答与讲解", "DeepSeek 生成 + 本地知识库 RAG 来源约束", "已验证"),
            ("路线推荐", "游玩时长、兴趣偏好、停留建议与推荐理由", "已实现"),
            ("数字人导览", "LiveTalking / Wav2Lip / EdgeTTS / WebRTC", "已跑通"),
            ("管理后台", "知识库、形象音色、反馈、行为数据看板", "已实现")]
    x0, y0 = 2.0, 3.45
    widths = [8.6, 15.4, 4.2]
    headers = ["赛题能力", "项目落点", "状态"]
    for i, h in enumerate(headers):
        rect(s, x0 + sum(widths[:i]), y0, widths[i], 0.9, ["blue", "teal", "green"][i], transparency=0)
        add_text(s, h, x0 + sum(widths[:i]) + 0.25, y0 + 0.22, widths[i] - 0.5, 0.35, 13.5, True, "bg", PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 1.0 + r * 1.8
        for i, val in enumerate(row):
            rect(s, x0 + sum(widths[:i]), y, widths[i], 1.55, "glass", "white", 24, False)
            add_text(s, val, x0 + sum(widths[:i]) + 0.34, y + 0.43, widths[i] - 0.68, 0.52, 15 if i != 1 else 14.3, True if i != 1 else False, "ink" if i != 1 else "text", PP_ALIGN.CENTER if i == 2 else None)
    add_text(s, "功能覆盖游客服务、AI 能力、数字人导览与景区后台运营。", 3.0, 14.25, 27.8, 0.7, 21, True, "ink", PP_ALIGN.CENTER)
    source_note(s, "项目源码与功能验收清单：游客端交互、RAG 问答、路线推荐、数字人导览、管理后台模块。")
    footer(s, 6)

    # 07 RAG
    s = blank_slide(prs)
    header(s, "RAG 知识库：降低幻觉，提升景区事实准确性", "先检索景区资料，再结合大模型生成可追溯回答", 7)
    stages = [("游客问题", "文本/语音输入", "blue"),
              ("知识检索", "约 50 个活跃文档", "teal"),
              ("模型生成", "DeepSeek deepseek-chat", "violet"),
              ("来源约束", "置信度与引用返回", "green")]
    for i, (h, b, c) in enumerate(stages):
        x = 2.0 + i * 7.55
        glass_panel(s, x, 4.0, 6.15, 4.6, 17)
        circle_label(s, str(i + 1), x + 2.2, 4.55, c, 1.5)
        add_text(s, h, x + 0.35, 6.55, 5.45, 0.45, 17, True, "ink", PP_ALIGN.CENTER)
        add_text(s, b, x + 0.35, 7.35, 5.45, 0.38, 12.5, color="muted", align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, "→", x + 6.15, 5.34, 1.1, 0.55, 27, True, "muted", PP_ALIGN.CENTER)
    big_metric(s, "≥90%", "事实准确率目标", 5.0, 11.25, 6.4, "blue", "赛题要求")
    big_metric(s, "0.90", "示例回答置信度", 13.7, 11.25, 6.4, "green", "灵山大佛问答")
    big_metric(s, "2.26s", "示例响应延迟", 22.4, 11.25, 6.4, "teal", "低于 5 秒目标")
    source_note(s, "赛题技术要求：事实准确率目标不低于 90%；项目本地知识库约 50 个活跃文档；示例 RAG 回答置信度 0.90、延迟约 2.26 秒。")
    footer(s, 7)

    # 08 Digital human
    s = blank_slide(prs)
    header(s, "数字人链路：从文本到可播报的实时音视频", "已修复本地连接失败问题，/offer 可返回 SDP 与 sessionid", 8)
    cover_picture(s, "lingling-guide-avatar.png", 1.6, 3.2, 7.7, 11.5)
    modules = [("讲解文本", "参考音色 / 景区讲解文案", "blue"),
               ("EdgeTTS", "文字转语音", "teal"),
               ("Wav2Lip", "口型驱动", "amber"),
               ("WebRTC", "浏览器实时播放", "green")]
    for i, (h, b, c) in enumerate(modules):
        x = 11.0 + i * 5.35
        glass_panel(s, x, 5.0, 4.55, 5.1, 18)
        circle_label(s, str(i + 1), x + 1.55, 5.58, c, 1.45)
        add_text(s, h, x + 0.28, 7.55, 4.0, 0.42, 15.5, True, "ink", PP_ALIGN.CENTER)
        add_text(s, b, x + 0.3, 8.35, 3.95, 0.65, 11.8, color="muted", align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, "→", x + 4.48, 6.25, 0.9, 0.55, 24, True, "muted", PP_ALIGN.CENTER)
    glass_panel(s, 11.2, 11.8, 20.2, 2.3, 20)
    add_text(s, "本地验证", 12.0, 12.15, 3.0, 0.42, 15, True, "green")
    add_text(s, "8010/index.html 可访问；/api/preset_voices、/api/admin/config、/api/avatar/tasks 和 /offer 均已验证。", 15.0, 12.08, 15.6, 0.72, 15, color="text")
    source_note(s, "本地数字人服务验证：8010/index.html、/api/preset_voices、/api/admin/config、/api/avatar/tasks、/offer 接口验证记录。")
    footer(s, 8)

    # 09 Route recommendation
    s = blank_slide(prs)
    header(s, "路线推荐样例：180 分钟历史文化路线", "把问答结果转化为可执行游览行动", 9)
    cover_picture(s, "scenic/lingshan-handdrawn-map.png", 1.35, 3.05, 15.2, 10.4)
    route = [("00-30", "九龙灌浴", "入口动线与文化开场"),
             ("30-80", "灵山大佛", "核心地标与主讲解"),
             ("80-125", "梵宫", "建筑与艺术体验"),
             ("125-180", "佛教文化博物馆", "收束与深度理解")]
    for i, (t, place, desc) in enumerate(route):
        y = 3.25 + i * 2.55
        glass_panel(s, 18.0, y, 12.8, 1.9, 18)
        add_text(s, t, 18.45, y + 0.36, 2.0, 0.38, 13.2, True, "blue", PP_ALIGN.CENTER)
        add_text(s, place, 21.1, y + 0.28, 4.2, 0.44, 16.5, True, "ink")
        add_text(s, desc, 21.1, y + 0.96, 8.8, 0.38, 12.8, color="muted")
    add_text(s, "路线推荐结合时间、兴趣偏好和停留建议，形成可执行游览决策。", 18.1, 14.2, 12.5, 0.72, 18, True, "ink")
    source_note(s, "路线推荐示例：输入 180 分钟游玩时长与历史文化偏好后生成的推荐路线。")
    footer(s, 9)

    # 10 Multi-end
    s = blank_slide(prs)
    header(s, "多端协同：游客端、数字人端、管理端一起工作", "按游客服务链路呈现系统协同关系", 10)
    cover_picture(s, "scenic/guide-map-clean.png", 2.0, 3.2, 9.6, 8.5)
    ends = [("游客端", "问答、讲解、路线推荐、移动入口", "blue"),
            ("数字人端", "头像、音色、参考音频、实时播报", "teal"),
            ("管理端", "知识库、反馈、行为数据、服务配置", "green")]
    for i, (h, b, c) in enumerate(ends):
        x = 13.2 + i * 6.35
        glass_panel(s, x, 4.25, 5.45, 7.5, 18)
        circle_label(s, str(i + 1), x + 1.9, 4.95, c, 1.55)
        add_text(s, h, x + 0.35, 7.1, 4.75, 0.45, 17.5, True, "ink", PP_ALIGN.CENTER)
        add_text(s, b, x + 0.45, 8.15, 4.55, 1.2, 13.3, color="muted", align=PP_ALIGN.CENTER)
        pill(s, "可演示", x + 1.55, 10.45, 2.35, c)
    add_text(s, "游客端、数字人端与管理端共同支撑完整景区导览服务。", 3.3, 13.8, 27.2, 0.72, 21, True, "ink", PP_ALIGN.CENTER)
    source_note(s, "项目本地页面与接口：游客端、数字人页面、管理后台功能入口。")
    footer(s, 10)

    # 11 Data dashboard
    s = blank_slide(prs)
    header(s, "运营数据证据：AI 导览可以沉淀景区经营洞察", "用真实数据看板证明项目具备行业适用性", 11)
    metrics = [("140447", "行为明细记录", "数据量支撑分析", "blue"),
               ("777", "匹配景区记录", "筛选后样本", "teal"),
               ("4.01h", "平均停留时长", "衡量游览深度", "green"),
               ("3.07", "平均满意度", "驱动反馈优化", "amber")]
    for i, (value, label, note, color) in enumerate(metrics):
        big_metric(s, value, label, 1.65 + i * 7.85, 3.8, 6.65, color, note)
    cover_picture(s, "scenic/photos/nine-dragons-bath.jpg", 2.0, 8.6, 11.2, 6.0)
    places = [("灵山大佛", 286, "blue"), ("拈花湾", 255, "teal"), ("灵山胜境", 236, "green")]
    max_v = 300
    add_text(s, "热门景点 Top 3", 15.0, 9.0, 6.0, 0.5, 18, True, "ink")
    for i, (name, val, c) in enumerate(places):
        y = 10.15 + i * 1.25
        add_text(s, name, 15.0, y, 3.2, 0.36, 13.5, True, "text")
        rect(s, 18.6, y + 0.06, 9.4, 0.34, "panel", transparency=0, radius=True)
        rect(s, 18.6, y + 0.06, 9.4 * val / max_v, 0.34, c, transparency=0, radius=True)
        add_text(s, str(val), 28.3, y - 0.04, 1.4, 0.36, 12.5, True, c)
    add_text(s, "数据范围：2025-01-03 至 2025-12-31", 15.0, 14.25, 12.0, 0.4, 11.5, color="muted")
    source_note(s, "景区行为数据看板与本地统计结果：行为明细 140447 条、匹配景区记录 777 条，数据范围 2025-01-03 至 2025-12-31。")
    footer(s, 11)

    # 12 Architecture
    s = blank_slide(prs)
    header(s, "系统架构：三服务解耦，AI 能力可替换扩展", "前端、业务 API、数字人服务分层运行，降低耦合和演示风险", 12)
    layers = [("前端体验层", "Vue 3 + Vite\n游客端 / 移动端 / 后台 / 数据可视化", "blue"),
              ("业务 API 层", "Python 后端\n问答 / 路线 / 知识库 / 反馈 / 统计", "teal"),
              ("AI 能力层", "DeepSeek + Qwen3-VL + RAG\n多模型协同生成", "violet"),
              ("数字人服务层", "LiveTalking / Wav2Lip / EdgeTTS / WebRTC\n实时音视频导览", "amber"),
              ("数据资源层", "SQLite / 本地文档 / 行为数据 / 景区图片", "green")]
    for i, (h, b, c) in enumerate(layers):
        y = 3.0 + i * 2.35
        glass_panel(s, 3.2, y, 27.2, 1.72, 18)
        add_text(s, h, 4.0, y + 0.42, 5.2, 0.45, 16.5, True, c)
        add_text(s, b, 10.2, y + 0.24, 18.8, 0.9, 14.2, color="text")
    add_text(s, "架构价值：三个本地服务已经跑通，便于现场演示与后续扩展。", 4.4, 15.4, 24.8, 0.65, 20, True, "ink", PP_ALIGN.CENTER)
    source_note(s, "项目架构与源码目录：Vue 3 + Vite 前端、Python 业务 API、LiveTalking/Wav2Lip/EdgeTTS/WebRTC 数字人服务。")
    footer(s, 12)

    # 13 Local validation
    s = blank_slide(prs)
    header(s, "本地运行验证：服务可启动、接口可调用、流程可演示", "用本地运行结果说明工程交付能力", 13)
    checks = [("8000", "后端 API", "问答、路线、看板通过", "green"),
              ("5173", "前端页面", "Vite 服务可访问", "blue"),
              ("8010", "数字人服务", "index、avatar、offer 通过", "teal"),
              ("55", "后端测试", "unittest 全部通过", "violet"),
              ("4", "前端用例", "Vitest 全部通过", "amber")]
    for i, (v, h, b, c) in enumerate(checks):
        x = 1.55 + i * 6.35
        glass_panel(s, x, 4.0, 5.45, 7.7, 18)
        add_text(s, v, x + 0.35, 4.75, 4.75, 0.9, 31, True, c, PP_ALIGN.CENTER)
        add_text(s, h, x + 0.45, 6.35, 4.55, 0.45, 16, True, "ink", PP_ALIGN.CENTER)
        add_text(s, b, x + 0.48, 7.45, 4.5, 0.9, 12.8, color="muted", align=PP_ALIGN.CENTER)
        pill(s, "PASS", x + 1.55, 10.15, 2.3, "green")
    add_text(s, "演示流程：游客端问答 → 路线推荐 → 数字人播报 → 后台看板。", 2.8, 14.3, 28.1, 0.72, 20.5, True, "ink", PP_ALIGN.CENTER)
    source_note(s, "本地运行与测试记录：后端 8000、前端 5173、数字人 8010；后端 unittest 55 个用例通过，前端 Vitest 4 个用例通过。")
    footer(s, 13)

    # 14 Demo script
    s = blank_slide(prs)
    header(s, "2 分钟展示流程：按游客服务链路演示", "每一步对应一个核心能力，保证展示节奏清晰", 14)
    script = [("0-20s", "打开游客端", "展示灵灵导游入口与多端导航", "blue"),
              ("20-50s", "RAG 问答", "询问“灵山大佛有什么特色？”并展示置信度/来源", "teal"),
              ("50-75s", "路线推荐", "输入 180 分钟 + 历史文化偏好", "green"),
              ("75-100s", "数字人播报", "连接 8010 页面，播报景区讲解文本", "amber"),
              ("100-120s", "后台看板", "展示知识库、反馈与行为数据", "violet")]
    for i, (tm, h, b, c) in enumerate(script):
        y = 3.25 + i * 2.4
        circle_label(s, str(i + 1), 2.0, y, c, 1.35)
        add_text(s, tm, 4.0, y + 0.08, 2.7, 0.42, 14.2, True, c)
        add_text(s, h, 7.15, y + 0.02, 5.5, 0.48, 18, True, "ink")
        add_text(s, b, 13.2, y + 0.08, 17.2, 0.45, 14.5, color="muted")
        if i < 4:
            arrow(s, 2.67, y + 1.35, 2.67, y + 2.25, "line")
    add_text(s, "每一步对应一个核心能力，展示路径清晰连贯。", 3.0, 15.6, 27.8, 0.65, 20, True, "ink", PP_ALIGN.CENTER)
    source_note(s, "演示流程依据：本地已跑通的游客端问答、路线推荐、数字人播报与后台看板功能。")
    footer(s, 14)

    # 15 Data sources
    s = blank_slide(prs)
    header(s, "数据来源与验证依据", "关键数字均来自赛题要求、本地运行测试和项目数据看板", 15)
    source_blocks = [
        ("赛题指标", [
            "评分结构：40% / 30% / 20% / 10%",
            "事实准确率目标：不低于 90%",
            "稳定响应延迟目标：不高于 5 秒",
        ], "blue"),
        ("本地问答验证", [
            "示例问题：灵山大佛有什么特色？",
            "RAG 置信度：0.90",
            "响应延迟：约 2.26 秒",
        ], "teal"),
        ("行为数据看板", [
            "行为明细记录：140447 条",
            "匹配景区记录：777 条",
            "数据范围：2025-01-03 至 2025-12-31",
            "平均停留时长：4.01 小时；平均满意度：3.07",
        ], "green"),
        ("工程验证记录", [
            "后端 API：127.0.0.1:8000",
            "前端页面：127.0.0.1:5173",
            "数字人服务：127.0.0.1:8010",
            "后端 unittest 55 个用例通过；前端 Vitest 4 个用例通过",
        ], "amber"),
    ]
    for i, (title_text, lines, color) in enumerate(source_blocks):
        x = 1.65 + (i % 2) * 15.7
        y = 3.4 + (i // 2) * 5.6
        glass_panel(s, x, y, 14.1, 4.65, 18)
        add_text(s, title_text, x + 0.55, y + 0.42, 6.0, 0.45, 17, True, color)
        add_lines(s, lines, x + 0.72, y + 1.25, 12.7, 2.8, 12.8, "text", gap=5, bullet=True)
    add_text(s, "说明：上述数据用于展示项目当前本地运行状态，最终提交材料可附部署说明、测试记录和数据看板截图作为支撑。", 2.2, 14.95, 29.0, 0.7, 15, True, "ink", PP_ALIGN.CENTER)
    footer(s, 15)

    # 16 Closing
    s = blank_slide(prs)
    cover_picture(s, "scenic/photos/bodhi-avenue.png", 22.2, 0, 11.7, H_CM)
    rect(s, 21.9, 0, 0.05, H_CM, "line", transparency=20)
    add_text(s, "可运行、可验证、可落地", 2.0, 4.2, 22.2, 1.05, 36, True, "ink")
    add_text(s, "灵灵导游把游客咨询、景点讲解、路线规划、数字人播报和景区运营闭环放进同一个 AI 导览系统。", 2.1, 5.85, 18.6, 1.25, 18, color="text")
    for i, (h, b, c) in enumerate([("功能闭环", "问、讲、游、管", "blue"),
                                   ("技术可信", "RAG + 多模型 + WebRTC", "teal"),
                                   ("工程交付", "本地三服务跑通", "green")]):
        x = 2.1 + i * 6.45
        glass_panel(s, x, 8.25, 5.6, 3.2, 18)
        add_text(s, h, x + 0.35, 8.78, 4.9, 0.45, 16, True, c, PP_ALIGN.CENTER)
        add_text(s, b, x + 0.35, 9.78, 4.9, 0.45, 13.2, color="text", align=PP_ALIGN.CENTER)
    add_text(s, "谢谢指导", 2.1, 15.65, 8.0, 0.62, 22, True, "ink")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    DOCS.mkdir(parents=True, exist_ok=True)
    out = make_deck()
    print(out)
